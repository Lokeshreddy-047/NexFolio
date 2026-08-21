import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG_LIGHT = RGBColor(248, 250, 252)
    CARD_BG = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(203, 213, 225)
    TEXT_MAIN = RGBColor(15, 23, 42)
    TEXT_MUTED = RGBColor(71, 85, 105)
    ACCENT_PURPLE = RGBColor(79, 70, 229)
    HEADER_DARK = RGBColor(30, 41, 59)

    def add_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_LIGHT
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = HEADER_DARK
        p.alignment = PP_ALIGN.CENTER

    def add_card(slide, left, top, width, height, title="", subtitle="", bg_color=CARD_BG, border_color=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)

        if title or subtitle:
            tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            if title:
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = ACCENT_PURPLE
                p.alignment = PP_ALIGN.CENTER
            if subtitle:
                p2 = tf.add_paragraph() if title else tf.paragraphs[0]
                p2.text = subtitle
                p2.font.size = Pt(10)
                p2.font.color.rgb = TEXT_MUTED
                p2.alignment = PP_ALIGN.CENTER
        return shape

    slide_layout = prs.slide_layouts[6]

    # SLIDE 1: Title Slide
    s1 = prs.slides.add_slide(slide_layout)
    add_background(s1)

    t_top = s1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.6))
    p = t_top.text_frame.paragraphs[0]
    p.text = "MAJOR PROJECT PRESENTATION"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = HEADER_DARK
    p.alignment = PP_ALIGN.CENTER

    t_main = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.933), Inches(2.2))
    tf = t_main.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "NexFolio: An Explainable AI (XAI) Framework for Intelligent Portfolio Risk Profiling, Real-Time Market Analytics, and Institutional Tax Optimization"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = HEADER_DARK
    p.alignment = PP_ALIGN.CENTER

    t_dept = s1.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(10.933), Inches(0.6))
    p = t_dept.text_frame.paragraphs[0]
    p.text = "Department of Computer Science and Engineering | CBIT"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

    t_pres = s1.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(5.5), Inches(1.4))
    tf = t_pres.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Presented By:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = HEADER_DARK
    p1 = tf.add_paragraph()
    p1.text = "Madupu Lokesh Reddy - 160122733047"
    p1.font.size = Pt(12)
    p1.font.color.rgb = TEXT_MAIN
    p2 = tf.add_paragraph()
    p2.text = "Patil Tejas                   - 160123733321"
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MAIN

    t_sup = s1.shapes.add_textbox(Inches(7.5), Inches(5.6), Inches(5.0), Inches(1.4))
    tf = t_sup.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Project Supervisor:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = HEADER_DARK
    p1 = tf.add_paragraph()
    p1.text = "Mr. Banothu Sai Kumar"
    p1.font.size = Pt(13)
    p1.font.color.rgb = TEXT_MAIN
    p2 = tf.add_paragraph()
    p2.text = "Assistant Professor, Dept. of CSE"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED

    # SLIDE 2: Title Justification
    s2 = prs.slides.add_slide(slide_layout)
    add_background(s2)
    add_header(s2, "TITLE JUSTIFICATION")

    cards2 = [
        ("“Explainable AI (XAI)”", "Replaces opaque 'black-box' predictions with game-theoretic TreeSHAP feature attributions. Translates complex mathematical weights into human-understandable risk drivers."),
        ("“Portfolio Risk Profiling”", "Multi-dimensional evaluation beyond simple P&L using 36 quantitative features across return momentum, annualized volatility, downside semi-variance, Sharpe/Sortino ratios, and market Beta."),
        ("“Real-Time Market Analytics”", "Decouples sub-millisecond in-memory valuation from heavy analytics using a Dual-Loop architecture and a 5-state Market Data Pedigree state machine (Live, Delayed, Reference)."),
        ("“Institutional Tax Optimization”", "Implements native statutory compliance with the new Income-tax Act, 2025 (Tax Year 2026–27), modeling STCG @ 20%, LTCG @ 12.5%, buyback taxation, and an 8-year loss carryforward bank.")
    ]
    coords2 = [
        (Inches(0.8), Inches(1.6), Inches(5.6), Inches(2.4)),
        (Inches(6.8), Inches(1.6), Inches(5.7), Inches(2.4)),
        (Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.4)),
        (Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.4))
    ]
    for (ctitle, cdesc), (l, t, w, h) in zip(cards2, coords2):
        add_card(s2, l, t, w, h, title=ctitle)
        tb = s2.shapes.add_textbox(l + Inches(0.2), t + Inches(0.7), w - Inches(0.4), h - Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cdesc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MAIN

    # SLIDE 3: Background
    s3 = prs.slides.add_slide(slide_layout)
    add_background(s3)
    add_header(s3, "BACKGROUND")

    tb3 = s3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "Quantitative Portfolio Management & AI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = HEADER_DARK
    p.space_after = Pt(14)

    bg_bullets = [
        "Modern Portfolio Theory (MPT) & Multi-Factor Models",
        "Exponential growth of retail investors (160M+ Demat accounts)",
        "Shift from simple P&L tracking to risk-adjusted metrics",
        "Supervised ML classification for portfolio risk tiers",
        "Statutory transition to Income-tax Act, 2025 (TY 2026–27)"
    ]
    for b in bg_bullets:
        p = tf3.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(10)

    img_path1 = "d:/nexfolio/fig1_architecture_flow_light.png"
    if os.path.exists(img_path1):
        s3.shapes.add_picture(img_path1, Inches(6.8), Inches(1.4), Inches(5.7), Inches(5.6))

    # SLIDE 4: Motivation
    s4 = prs.slides.add_slide(slide_layout)
    add_background(s4)
    add_header(s4, "MOTIVATION")

    mots = [
        "Retail Capital Exposure",
        "Black-Box AI Distrust",
        "Budget 2026-27 Tax Overhaul",
        "Cross-Sector Market Volatility",
        "High-Latency System Freezes",
        "Audit & Regulatory Filings"
    ]
    coords_m = [
        (Inches(1.2), Inches(1.8), Inches(5.2), Inches(1.4)),
        (Inches(6.9), Inches(1.8), Inches(5.2), Inches(1.4)),
        (Inches(1.2), Inches(3.5), Inches(5.2), Inches(1.4)),
        (Inches(6.9), Inches(3.5), Inches(5.2), Inches(1.4)),
        (Inches(1.2), Inches(5.2), Inches(5.2), Inches(1.4)),
        (Inches(6.9), Inches(5.2), Inches(5.2), Inches(1.4))
    ]
    for m_text, (l, t, w, h) in zip(mots, coords_m):
        add_card(s4, l, t, w, h)
        tb = s4.shapes.add_textbox(l + Inches(0.2), t + Inches(0.4), w - Inches(0.4), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = m_text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = HEADER_DARK
        p.alignment = PP_ALIGN.CENTER

    # SLIDE 5: The Core Challenge
    s5 = prs.slides.add_slide(slide_layout)
    add_background(s5)
    add_header(s5, "THE PORTFOLIO RISK & BLACK-BOX CHALLENGE")

    tb5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.2), Inches(5.0))
    tf5 = tb5.text_frame
    tf5.word_wrap = True
    chal_bullets = [
        "Model Opacity in deep neural networks",
        "High feature multi-collinearity & noise",
        "Inference latency bottlenecks (>40ms)",
        "Outdated 365-day tax approximations"
    ]
    for b in chal_bullets:
        p = tf5.add_paragraph() if tf5.paragraphs[0].text else tf5.paragraphs[0]
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(14)

    img_path2 = "d:/nexfolio/fig2_benchmark_tradeoff_light.png"
    if os.path.exists(img_path2):
        s5.shapes.add_picture(img_path2, Inches(6.2), Inches(1.4), Inches(6.3), Inches(5.4))

    # SLIDE 6: Literature Survey (Accurate 4 Papers from PDFs)
    s6 = prs.slides.add_slide(slide_layout)
    add_background(s6)
    add_header(s6, "LITERATURE SURVEY")

    rows, cols = 5, 5
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4)
    table_shape = s6.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.3)
    table.columns[3].width = Inches(2.3)
    table.columns[4].width = Inches(2.433)

    headers = ["Authors", "Title", "Methods", "Advantages", "Disadvantages"]
    for c_idx, h_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = HEADER_DARK

    papers_data = [
        ("P. Singh et al.\n(IEEE Access, 2023)\n[Base Paper]", "Harnessing a Hybrid CNN-LSTM Model for Portfolio Performance", "1D-CNN, LSTM, Markowitz Mean-Variance", "Captures temporal trends; multi-scale feature learning", "Opaque black-box; high compute latency (~45ms); no tax math."),
        ("I. Aruleba, Y. Sun\n(IEEE Access, 2024)", "Effective Credit Risk Prediction Using Ensemble Classifiers", "Random Forest, XGBoost, SMOTE-ENN, TreeSHAP", "Ensemble superiority on tabular data; exact Shapley values", "Binary credit default focus; lacks multi-class portfolio tiers."),
        ("I. Aruleba, Y. Sun\n(IEEE Access, 2025)", "An Improved Ensemble Method With Data Resampling", "Stacked Ensemble (RF, LR, CNN + MLP), SMOTE-ENN", "Mitigates extreme class skew; reduces individual model variance", "High stacking overhead; static batch testing without streaming."),
        ("Vijayanand, Smrithy\n(IDT Journal, 2025)", "Explainable AI-Enhanced Ensemble Learning for Financial Fraud", "Voting Ensemble (XGB, RF, DT) + TreeSHAP", "Achieves 99.9% accuracy; provides feature-level audit trust", "Limited to transaction fraud; lacks portfolio metrics or tax rules.")
    ]

    for r_idx, row_data in enumerate(papers_data, start=1):
        for c_idx, cell_value in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = cell_value
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_MAIN

    # SLIDE 7: Research Gap
    s7 = prs.slides.add_slide(slide_layout)
    add_background(s7)
    add_header(s7, "RESEARCH GAP")

    gap_data = [
        ("Black-Box Risk Models", "TreeSHAP Attribution Layer", "Sub-second game-theoretic risk drivers for every prediction"),
        ("Feature Collinearity", "Dual L1 + L2 Regularization", "Prevents overfitting on correlated ratios, achieving 97% accuracy"),
        ("Analytical Latency", "Dual-Loop Engine (<1ms SSE)", "Decouples live quote streaming from heavy persistence"),
        ("Outdated Tax Engines", "Income-tax Act, 2025 Suite", "Exact calendar-month rules, buybacks, and 8-year loss bank")
    ]

    for idx, (g1, g2, g3) in enumerate(gap_data):
        top_pos = Inches(1.6 + idx * 1.3)
        b1 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(3.2), Inches(1.0))
        b1.fill.solid()
        b1.fill.fore_color.rgb = RGBColor(241, 245, 249)
        b1.line.color.rgb = BORDER_COLOR
        p = b1.text_frame.paragraphs[0]
        p.text = g1
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = HEADER_DARK
        p.alignment = PP_ALIGN.CENTER

        t_arr1 = s7.shapes.add_textbox(Inches(4.1), top_pos + Inches(0.2), Inches(0.5), Inches(0.5))
        p = t_arr1.text_frame.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_PURPLE

        b2 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), top_pos, Inches(3.4), Inches(1.0))
        b2.fill.solid()
        b2.fill.fore_color.rgb = RGBColor(238, 242, 255)
        b2.line.color.rgb = ACCENT_PURPLE
        p = b2.text_frame.paragraphs[0]
        p.text = g2
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_PURPLE
        p.alignment = PP_ALIGN.CENTER

        t_arr2 = s7.shapes.add_textbox(Inches(8.2), top_pos + Inches(0.2), Inches(0.5), Inches(0.5))
        p = t_arr2.text_frame.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_PURPLE

        tb3 = s7.shapes.add_textbox(Inches(8.8), top_pos + Inches(0.1), Inches(3.7), Inches(0.8))
        tf = tb3.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = g3
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_MUTED

    # SLIDE 8: Single Concise Problem Statement Slide (Two Clean Cards)
    s8 = prs.slides.add_slide(slide_layout)
    add_background(s8)
    add_header(s8, "PROBLEM STATEMENT")

    # Card 1: ⚠️ THE PROBLEM
    add_card(s8, Inches(1.2), Inches(1.6), Inches(10.933), Inches(2.4), bg_color=RGBColor(254, 242, 242), border_color=RGBColor(239, 68, 68))
    tb1 = s8.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(2.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "⚠️ THE CURRENT CHALLENGE:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(185, 28, 28)
    p.space_after = Pt(6)

    prob_items = [
        "Model Opacity: Black-box deep neural networks obscure underlying risk drivers.",
        "Analytical Latency: Heavy compute bottlenecks (>40ms) stall live ticker streams.",
        "Statutory Ignorance: Outdated tools ignore tax drag under the new Income-tax Act, 2025."
    ]
    for pi in prob_items:
        p = tf1.add_paragraph()
        p.text = "• " + pi
        p.font.size = Pt(11.5)
        p.font.color.rgb = HEADER_DARK

    # Card 2: 💡 PROPOSED OBJECTIVE
    add_card(s8, Inches(1.2), Inches(4.3), Inches(10.933), Inches(2.6), bg_color=RGBColor(240, 253, 250), border_color=RGBColor(13, 148, 136))
    tb2 = s8.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.333), Inches(2.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "💡 OUR PROPOSED OBJECTIVE (NexFolio):"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 118, 110)
    p.space_after = Pt(6)

    sol_items = [
        "Explainable ML: Regularized XGBoost (97% Acc) with game-theoretic TreeSHAP attributions.",
        "Dual-Loop Engine: Decouple <1ms live quote streaming from background deep analytics.",
        "Institutional Tax Suite: Native STCG (20%), LTCG (12.5%), and 8-year Loss Harvesting Bank."
    ]
    for si in sol_items:
        p = tf2.add_paragraph()
        p.text = "• " + si
        p.font.size = Pt(11.5)
        p.font.color.rgb = HEADER_DARK

    # SLIDE 9: Objectives (6 Numbered Cards)
    s9 = prs.slides.add_slide(slide_layout)
    add_background(s9)
    add_header(s9, "OBJECTIVES")

    objs = [
        ("01", "36-Feature Pipeline", "Capture return momentum, volatility, downside risk, Beta, and 18 sector weights."),
        ("02", "XGBoost & TreeSHAP", "Train regularized gradient-boosted trees (>95% accuracy) with sub-second SHAP attributions."),
        ("03", "4-Pillar Scorecard", "Formulate an inspectable 0–100 score across Diversification, Volatility, Efficiency, and Drawdown."),
        ("04", "What-If Sandbox", "Enable in-memory trade simulations with instant delta calculations without database writes."),
        ("05", "Dual-Loop Engine", "Decouple <1ms live quote streaming from slow persistence with a 5-state Pedigree FSM."),
        ("06", "Statutory Tax Suite", "Implement Income-tax Act, 2025 STCG @ 20%, LTCG @ 12.5%, Buybacks, and 8-Year Loss Bank.")
    ]

    coords_o = [
        (Inches(0.8), Inches(1.6), Inches(5.6), Inches(1.6)),
        (Inches(0.8), Inches(3.4), Inches(5.6), Inches(1.6)),
        (Inches(0.8), Inches(5.2), Inches(5.6), Inches(1.6)),
        (Inches(6.8), Inches(1.6), Inches(5.7), Inches(1.6)),
        (Inches(6.8), Inches(3.4), Inches(5.7), Inches(1.6)),
        (Inches(6.8), Inches(5.2), Inches(5.7), Inches(1.6))
    ]

    for (num, otitle, odesc), (l, t, w, h) in zip(objs, coords_o):
        add_card(s9, l, t, w, h)
        nb = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l + Inches(0.15), t + Inches(0.15), Inches(0.6), Inches(0.6))
        nb.fill.solid()
        nb.fill.fore_color.rgb = ACCENT_PURPLE
        p_n = nb.text_frame.paragraphs[0]
        p_n.text = num
        p_n.font.size = Pt(12)
        p_n.font.bold = True
        p_n.font.color.rgb = RGBColor(255, 255, 255)
        p_n.alignment = PP_ALIGN.CENTER

        tb = s9.shapes.add_textbox(l + Inches(0.85), t + Inches(0.15), w - Inches(0.95), h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = otitle
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = HEADER_DARK
        p_sub = tf.add_paragraph()
        p_sub.text = odesc
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = TEXT_MUTED

    # SLIDE 10: Conclusion
    s10 = prs.slides.add_slide(slide_layout)
    add_background(s10)
    add_header(s10, "CONCLUSION")

    tb10 = s10.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(5.2))
    tf10 = tb10.text_frame
    tf10.word_wrap = True

    c_bullets = [
        "Investigate black-box limitations in portfolio risk models and formulate an explainable 36-feature quantitative pipeline.",
        "Develop an institutional ML architecture combining regularized XGBoost with TreeSHAP for sub-second risk attributions.",
        "Formulate a deterministic 4-pillar health scorecard (0–100) and an in-memory What-If trade simulation sandbox.",
        "Implement a Dual-Loop valuation architecture to decouple <1ms live quote streaming from background analytics.",
        "Implement statutory tax compliance for the Income-tax Act, 2025 featuring STCG 20%, LTCG 12.5%, and an 8-year Loss Bank."
    ]

    for i, cb in enumerate(c_bullets):
        p = tf10.paragraphs[0] if i == 0 else tf10.add_paragraph()
        p.text = "• " + cb
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(14)

    # SLIDE 11: References
    s11 = prs.slides.add_slide(slide_layout)
    add_background(s11)
    add_header(s11, "REFERENCES")

    tb11 = s11.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
    tf11 = tb11.text_frame
    tf11.word_wrap = True

    refs = [
        "[1] P. Singh, M. Jha, M. Sharaf, M. A. El-Meligy, and T. R. Gadekallu, 'Harnessing a Hybrid CNN-LSTM Model for Portfolio Performance: A Case Study on Stock Selection and Optimization,' IEEE Access, vol. 11, pp. 45120–45132, Sep. 2023.",
        "[2] I. Aruleba and Y. Sun, 'Effective Credit Risk Prediction Using Ensemble Classifiers With Model Explanation,' IEEE Access, vol. 12, pp. 115015–115025, Aug. 2024.",
        "[3] I. Aruleba and Y. Sun, 'An Improved Ensemble Method With Data Resampling for Credit Risk Prediction,' IEEE Access, vol. 13, pp. 128940–128952, Apr. 2025.",
        "[4] D. Vijayanand and G. S. Smrithy, 'Explainable AI-enhanced ensemble learning for financial fraud detection in mobile money transactions,' Intelligent Decision Technologies, vol. 19, no. 1, pp. 52–67, 2025.",
        "[5] S. M. Lundberg and S.-I. Lee, 'A Unified Approach to Interpreting Model Predictions,' in Advances in Neural Information Processing Systems (NeurIPS), 2017.",
        "[6] Ministry of Finance, Government of India, 'The Income-tax Act, 2025 and Union Budget 2026–27 Statutory Provisions,' New Delhi, 2026."
    ]

    for i, rf in enumerate(refs):
        p = tf11.paragraphs[0] if i == 0 else tf11.add_paragraph()
        p.text = rf
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(8)

    # SLIDE 12: Thank You
    s12 = prs.slides.add_slide(slide_layout)
    add_background(s12)

    tb12 = s12.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.733), Inches(1.5))
    p = tb12.text_frame.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = HEADER_DARK
    p.alignment = PP_ALIGN.CENTER

    out_file = "d:/nexfolio/NexFolio_Review1_Presentation_Deck.pptx"
    prs.save(out_file)
    print(f"Presentation deck saved successfully to: {out_file}")

if __name__ == "__main__":
    create_presentation()
